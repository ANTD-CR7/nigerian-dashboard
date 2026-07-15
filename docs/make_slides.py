"""Build a 21-slide NPEDATA defence deck (16:9) in the platform's palette.
Run:  python docs/make_slides.py   ->  docs/NPEDATA_Defense_Slides.pptx
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIG = Path(__file__).parent / "figures"
OUT = Path(__file__).parent / "NPEDATA_Defense_Slides.pptx"

# ── NPEDATA palette ──
BLACK   = RGBColor(0x07, 0x08, 0x0D)
PANEL   = RGBColor(0x12, 0x14, 0x1F)
CREAM   = RGBColor(0xF5, 0xF0, 0xE8)
GREEN   = RGBColor(0x00, 0xA3, 0x62)
GOLD    = RGBColor(0xF4, 0xA0, 0x17)
SLATE   = RGBColor(0x9A, 0x96, 0x8F)
HEAD_F  = "Georgia"        # serif proxy for the site's Playfair headings
BODY_F  = "Segoe UI"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def bg(slide, color=BLACK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.text_frame.word_wrap = True
    return tb


def set_run(r, text, size, color, font=BODY_F, bold=False, italic=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic


def accent_bar(slide, l=Inches(0.9), t=Inches(0.62), w=Inches(0.55)):
    sh = slide.shapes.add_shape(1, l, t, w, Inches(0.06))  # rectangle
    sh.fill.solid(); sh.fill.fore_color.rgb = GOLD
    sh.line.fill.background()


def header(slide, eyebrow, title, n):
    accent_bar(slide)
    # eyebrow
    e = box(slide, Inches(0.9), Inches(0.72), Inches(9.5), Inches(0.4))
    p = e.text_frame.paragraphs[0]
    set_run(p.add_run(), eyebrow.upper(), 12, GREEN, BODY_F, bold=True)
    p.runs[0].font.spacing = Pt(2) if hasattr(p.runs[0].font, "spacing") else None
    # title
    tt = box(slide, Inches(0.88), Inches(1.05), Inches(11.4), Inches(1.2))
    tp = tt.text_frame.paragraphs[0]
    set_run(tp.add_run(), title, 32, CREAM, HEAD_F, italic=True)
    # slide number chip
    num = box(slide, Inches(12.35), Inches(6.95), Inches(0.8), Inches(0.4))
    np_ = num.text_frame.paragraphs[0]; np_.alignment = PP_ALIGN.RIGHT
    set_run(np_.add_run(), f"{n:02d} / 21", 10, SLATE, BODY_F)


def bullets(slide, items, top=Inches(2.35), left=Inches(1.0), width=Inches(11.2),
            size=18, gap=10):
    tb = box(slide, left, top, width, Inches(4.4))
    tf = tb.text_frame
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        set_run(p.add_run(), "▸  ", size, GOLD, BODY_F, bold=True)
        if lead:
            set_run(p.add_run(), lead, size, CREAM, BODY_F, bold=True)
        if rest:
            set_run(p.add_run(), (" — " if lead else "") + rest, size, SLATE, BODY_F)
    return tb


def add_image(slide, path, l, t, max_w, max_h):
    from PIL import Image
    iw, ih = Image.open(path).size
    ratio = min(max_w / iw, max_h / ih)
    w, h = int(iw * ratio), int(ih * ratio)
    # center within the box
    pic = slide.shapes.add_picture(str(path), l + (max_w - w) // 2, t + (max_h - h) // 2, w, h)
    return pic


def notes(slide, text):
    """Attach speaker notes — the presenter's context, never shown on screen."""
    if text:
        slide.notes_slide.notes_text_frame.text = text


def content(eyebrow, title, n, items=None, img=None, img_side="right", caption=None,
            sub=None, note=None):
    s = prs.slides.add_slide(BLANK); bg(s)
    header(s, eyebrow, title, n)
    b_text_top, b_img_top, img_top = Inches(2.7), Inches(2.55), Inches(2.35)
    if sub:
        sb = box(s, Inches(0.92), Inches(1.92), Inches(11.6), Inches(0.7))
        set_run(sb.text_frame.paragraphs[0].add_run(), sub, 15, GOLD, BODY_F, italic=True)
        b_text_top, b_img_top, img_top = Inches(3.15), Inches(3.05), Inches(2.85)
    if img and items:
        # split: bullets left, image right
        bullets(s, items, top=b_img_top, left=Inches(1.0), width=Inches(6.0), size=18, gap=13)
        add_image(s, img, Inches(7.25), img_top, Inches(5.4), Inches(4.0))
    elif img:
        add_image(s, img, Inches(1.4), img_top, Inches(10.5), Inches(4.2))
    elif items:
        bullets(s, items, top=b_text_top, size=22, gap=18)
    if caption:
        c = box(s, Inches(1.0), Inches(6.9), Inches(11.0), Inches(0.4))
        cp = c.text_frame.paragraphs[0]
        set_run(cp.add_run(), caption, 11, SLATE, BODY_F, italic=True)
    notes(s, note)
    return s

# ─────────────────────────── SLIDE 1: TITLE ───────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, Inches(0.9), Inches(1.0), Inches(0.7))
eb = box(s, Inches(0.9), Inches(1.15), Inches(11), Inches(0.4))
set_run(eb.text_frame.paragraphs[0].add_run(),
        "CALEB UNIVERSITY, LAGOS  ·  B.Sc. COMPUTER SCIENCE  ·  FINAL YEAR PROJECT", 12, GREEN, BODY_F, bold=True)
tt = box(s, Inches(0.85), Inches(1.7), Inches(11.6), Inches(2.6))
para = tt.text_frame.paragraphs[0]
set_run(para.add_run(),
        "Design and Development of a Nigerian Public Economic Data Aggregation and Analytics Platform with Open API Access",
        34, CREAM, HEAD_F, italic=True)
sub = box(s, Inches(0.9), Inches(4.15), Inches(11), Inches(0.5))
set_run(sub.text_frame.paragraphs[0].add_run(), "A 2020–2026 Case Study  ·  “NPEDATA”", 18, GOLD, HEAD_F, italic=True)
meta = box(s, Inches(0.9), Inches(5.0), Inches(11.5), Inches(2.0))
lines = [
    ("By:  ", "Taoheed Abdulmanan Olaosebikan   (22/10267)"),
    ("Supervisor:  ", "Miss Ilori Deborah"),
    ("Head of Department:  ", "Dr. Ayorinde Oduroye"),
    ("Department:  ", "Computer Science, College of Science and Information Science (COSIS)"),
    ("Date:  ", "July 2026"),
]
tf = meta.text_frame
for i, (k, v) in enumerate(lines):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.space_after = Pt(5)
    set_run(p.add_run(), k, 14, GREEN, BODY_F, bold=True)
    set_run(p.add_run(), v, 14, CREAM, BODY_F)
notes(s, "Greet the panel, give my name and the project title in one sentence: a platform that "
         "pulls Nigeria's scattered public economic data into one place and serves it to both "
         "people and programs. Keep this to about twenty seconds, then move on.")

# ─────────────────────────── SLIDE 2 ───────────────────────────
content("Introduction", "Background of the Study", 2,
    sub="Reliable data exists, but its published form makes it hard to use.",
    items=[
        ("Data drives decisions", "government, researchers, journalists and the public all rely on economic figures."),
        ("Nigeria's sources are credible", "the figures come from the CBN, the NBS and the World Bank."),
        ("But the data is scattered", "it is spread across many websites as PDFs and spreadsheets."),
        ("So it is hard to use", "any analysis must begin by collecting and cleaning the data by hand."),
    ],
    note="Open here. The raw material is not the problem: the CBN, NBS and World Bank produce "
         "reliable figures. The problem is the form, PDFs and inconsistent spreadsheets, so anyone "
         "who wants to analyse the data must first collect and clean it by hand before any real "
         "work begins. Removing that effort is what this project is about.")

# ─────────────────────────── SLIDE 3 ───────────────────────────
content("The Problem", "Statement of the Problem", 3,
    sub="Four barriers make Nigeria's public data hard to use in practice.",
    items=[
        ("Not machine-readable", "the data is locked inside PDF documents."),
        ("No common structure", "units and time periods differ from one source to the next."),
        ("No open API", "there is no free way to query the data with a program."),
        ("Repeated effort", "every user must redo the same collection and cleaning."),
    ],
    note="These four are the specific problems I target, and each maps to a later design "
         "decision: not machine-readable maps to the pipeline; no common structure to the tidy "
         "database model; no open API to the FastAPI HATEOAS API; repeated effort to one shared, "
         "reusable store. Explain each briefly, do not just read them.")

# ─────────────────────────── SLIDE 4 ───────────────────────────
s = content("Objectives", "Aim and Objectives", 4,
    note="State the aim in one sentence, then note that the six objectives map directly onto the "
         "architecture I am about to show: collect, standardise, analyse, API, dashboard, "
         "evaluate. If asked, each objective was met, and the evaluation objective is the "
         "testing chapter.")
a = box(s, Inches(1.0), Inches(2.25), Inches(11.3), Inches(1.1))
ap = a.text_frame.paragraphs[0]
set_run(ap.add_run(), "Aim:  ", 18, GOLD, BODY_F, bold=True)
set_run(ap.add_run(), "To build a platform that gathers Nigeria's public economic data into one standard store, available through a dashboard and a free, open API.", 18, CREAM, BODY_F)
objs = [
    "Collect economic indicators from the CBN, NBS and World Bank into one repository.",
    "Standardise them into one relational database model.",
    "Provide analytics: change, trend and correlation.",
    "Build a free REST API with HATEOAS and no login.",
    "Build a clear, interactive dashboard.",
    "Test the system for correctness, usability and accessibility.",
]
tb = box(s, Inches(1.0), Inches(3.6), Inches(11.3), Inches(3.3)); tf = tb.text_frame
for i, v in enumerate(objs):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(8)
    set_run(p.add_run(), f"{i+1}.   ", 16, GREEN, BODY_F, bold=True)
    set_run(p.add_run(), v, 16, CREAM, BODY_F)

# ─────────────────────────── SLIDE 5 ───────────────────────────
content("Scope", "Scope of the Project", 5,
    sub="A clear boundary: what the project covers, and what it does not.",
    items=[
        ("It is a prototype", "a working system with 122 indicators and about 12,100 records."),
        ("It is not a national system", "not official infrastructure, not a bank, not real-time."),
        ("It uses no AI", "the analytics are kept simple and transparent on purpose."),
        ("Data is collected manually", "the figures reflect the latest snapshot, not a live feed."),
    ],
    note="This is my honesty slide, and I say it plainly: a working prototype and a controlled "
         "case study, not national infrastructure, not a bank, not real-time, and no AI, by "
         "choice, because the goal is transparency and trust rather than prediction. Saying this "
         "early answers the is-this-a-national-system question before it is asked.")

# ─────────────────────────── SLIDE 6 ───────────────────────────
content("Literature Review", "Key Concepts", 6,
    sub="The project applies established, well-known principles rather than inventing new ones.",
    items=[
        ("Open data", "information is only useful when it is machine-readable."),
        ("Tidy data", "one record per row, with metadata kept separate (Wickham, 2014)."),
        ("REST and HATEOAS", "an API that describes itself to its users (Fielding, 2000)."),
        ("Honest charts", "avoid misleading tricks such as dual axes (Tufte, 1983)."),
    ],
    note="If challenged on any design decision, I can name its source: tidy data from Wickham, "
         "REST and HATEOAS from Fielding, honest visualisation from Tufte and Anscombe, and the "
         "open-data principles. The project is applied engineering on an established base, not "
         "invented from scratch, which is a strength.")

# ─────────────────────────── SLIDE 7 ───────────────────────────
content("Existing Systems", "Related Systems and the Gap", 7,
    sub="Good systems exist, but none meet all three needs at once.",
    items=[
        ("FRED and the World Bank", "excellent access, but little detailed Nigerian data."),
        ("Trading Economics, Statista", "mostly paywalled, and resell the same figures."),
        ("The CBN and NBS", "reliable sources, but they offer no API."),
        ("The gap", "no system offers Nigerian data, open access and honesty together."),
    ],
    note="I am not claiming these systems are bad, FRED and the World Bank are excellent. I am "
         "pointing to one empty space: detailed Nigerian data, free machine access, and "
         "statistical-honesty features, together. That combination is what is missing, and it is "
         "what I built. Frame it as a gap, not a criticism.")

# ─────────────────────────── SLIDE 8 ───────────────────────────
content("Motivation", "Why This Has Not Been Built", 8,
    sub="The obstacle is mainly institutional, not technical (all points are cited).",
    items=[
        ("A framework already exists", "a government interoperability framework since 2018."),
        ("Data is treated as an asset", "agencies are often reluctant to share it (Eleanya, 2026)."),
        ("Rules are not enforced", "153 of 250 agencies fell below the FOI benchmark (Ogunyale & Osho, 2023)."),
        ("This project shows it is doable", "the technical barrier is low."),
    ],
    note="Keep this evidence-based and attributed, and do not editorialise about government. I am "
         "reporting published work: the interoperability framework exists since 2018, an audit "
         "found 153 of 250 agencies below the FOI benchmark, and the 2023 Data Protection Act "
         "set a precedent for enforcement. My point is narrow: the barrier is not engineering, "
         "which a single-student build in one year demonstrates. If a panellist raises politics, "
         "return to that narrow point.")

# ─────────────────────────── SLIDE 9 ───────────────────────────
content("Methodology", "Research Methodology", 9,
    sub="An iterative approach suited a solo project with requirements that emerged over time.",
    items=[
        ("Built in stages", "data model first, then API, then dashboard, then analytics."),
        ("Checked every stage", "displayed figures were verified against the database each cycle."),
        ("Not Waterfall or Scrum", "requirements were not fully known, and this was a one-person project."),
    ],
    note="Justify the methodology if asked: iterative and incremental, because I could not know "
         "every source quirk up front and I am one developer, so Waterfall and full team Scrum "
         "both fit poorly. The discipline that mattered was verifying every displayed figure "
         "against the database each cycle, which is why design and testing are closely linked.")

# ─────────────────────────── SLIDE 10 ───────────────────────────
content("System Design", "System Architecture", 10, img=FIG / "fig3_1_architecture.png",
    sub="A seven-stage pipeline transforms the data from source to screen.",
    caption="Collect → Validate → Standardise → Store → Analyse → API → Present.",
    note="This diagram is the heart of the project, walk it left to right: collect, validate, "
         "standardise, store, analyse, then present it two ways, the API and the dashboard. Spend "
         "time on validate and standardise, because that is where the value is added and what "
         "separates this from simply re-displaying the source websites.")

# ─────────────────────────── SLIDE 11 ───────────────────────────
content("System Design", "Database Design", 11, img=FIG / "fig3_4_erd.png",
    sub="One tidy table can hold series of every time frequency.",
    items=[
        ("Three linked tables", "sources, indicators and observations."),
        ("One observations table", "stores every series as (indicator, date, value)."),
        ("Built-in integrity", "keys and a uniqueness rule prevent duplicate records."),
        ("Units are stored", "each series keeps its unit, avoiding thousands-vs-millions errors."),
    ], caption="Figure 3.4 — Entity-relationship diagram.",
    note="The design idea is separating the value from its metadata. One observations table, in "
         "the shape indicator-date-value, lets any frequency live together and be queried by one "
         "engine. The uniqueness rule makes duplicate records impossible, and storing each "
         "series' unit prevents the thousands-versus-millions error I caught and fixed during "
         "the build.")

# ─────────────────────────── SLIDE 12 ───────────────────────────
content("Implementation", "The Open API", 12, img=FIG / "fig4_11_hateoas_explorer.png",
    sub="A free, self-describing API — the machine access the original sources lack.",
    items=[
        ("Built with FastAPI", "15 data endpoints, and no login is required."),
        ("Self-describing", "every response includes links to related data (HATEOAS Level 3)."),
        ("Fully navigable", "the whole API can be explored from one starting point."),
        ("Ready for machines", "it also serves AI agents through an MCP server."),
    ], caption="Figure 4.11 — HATEOAS Explorer browsing the live API.",
    note="HATEOAS Level 3 means every response includes links to related data, so the whole API "
         "can be explored from the start page with no manual, and I can prove it live with the "
         "Explorer. This is what makes it a platform others can build on, not just a website. "
         "Very few APIs reach Level 3, so this is a real strength.")

# ─────────────────────────── SLIDE 13 ───────────────────────────
content("Implementation", "The Dashboard", 13, img=FIG / "fig4_4_analyst_dial.png",
    sub="One clear interface that serves both the public and researchers.",
    items=[
        ("Each page tells a story", "what happened, why it matters, and how to read it."),
        ("A Reader / Analyst switch", "the same page serves the public and researchers."),
        ("Honest charts", "no misleading dual axes, and units are always shown."),
        ("Fully accessible", "it scores 100 out of 100 for accessibility."),
    ], caption="Figure 4.4 — The Reader/Analyst dial revealing statistical detail.",
    note="The Reader/Analyst switch is the interface idea: one page, one codebase, two levels of "
         "depth, so I never maintain two separate interfaces. Charts are honest throughout, no "
         "dual axes and units always shown, with a verified 100/100 accessibility score.")

# ─────────────────────────── SLIDE 14 — ANALYTICS ───────────────────────────
content("Analytics", "The Analytics Engine", 14, img=FIG / "fig4_9_compare_significance.png",
    sub="The platform measures relationships in the data, not just displays numbers.",
    items=[
        ("A full profile for any indicator", "latest value, change, range, volatility and trend."),
        ("Correlation with significance", "shown with an R² value and a p-value."),
        ("A correlation matrix", "compares 12 key indicators against each other at once."),
        ("Transparent methods", "all classical statistics, no hidden AI model."),
    ], caption="Figure 4.9 — Correlation reported with R² and a significance p-value.",
    note="This answers the does-it-actually-do-analytics question directly. Any of the 122 "
         "indicators gets a full live profile; correlations come with an R-squared value and a "
         "p-value; and the matrix compares 12 key series at once. Using classical, transparent "
         "statistics with no hidden model is a deliberate choice, and I should say so clearly.")

# ─────────────────────────── SLIDE 15 ───────────────────────────
content("Feature", "Reform Impact", 15, img=FIG / "fig4_15_reform_impact.png",
    sub="The data applied to a real national debate, without taking a side.",
    items=[
        ("Before and after June 2023", "it compares every indicator across the reforms."),
        ("Computed from the data", "the figures are calculated live from the stored records."),
        ("It takes no side", "critics and supporters see the very same numbers."),
    ], caption="Figure 4.15 — Reform Impact: before and after June 2023, neutral readings.",
    note="Handle this one calmly and stay neutral. The platform compares every main indicator "
         "before and after June 2023 and lets a critic and a supporter both cite real numbers. If "
         "a panellist pushes a political view, I do not argue it; I show that the neutrality is "
         "the point, the tool makes the data checkable, it does not settle the argument.")

# ─────────────────────────── SLIDE 16 ───────────────────────────
content("Integrity", "Honesty and Validation", 16, img=FIG / "fig4_12_playground.png",
    sub="The platform flags results that could mislead the reader.",
    items=[
        ("Warns about false links", "it flags a correlation caused only by a shared trend."),
        ("Separates two ideas", "statistical significance is kept apart from strength."),
        ("Public data-check tool", "anyone can validate a file, and it can never change the data."),
    ], caption="Figure 4.12 — Pipeline Playground: per-row verdicts, nothing written.",
    note="Lead with the honesty story, it is my strongest point. The guard compares the ordinary "
         "correlation with the de-trended one, for example CBN assets versus currency falls from "
         "0.81 to 0.03 and is flagged. Significance is kept separate from strength, and the "
         "validation tool is public but can never write to the database.")

# ─────────────────────────── SLIDE 17 ───────────────────────────
content("Testing", "Testing and Validation", 17,
    sub="Correctness was tested and shown, not simply assumed.",
    items=[
        ("24 automated tests", "they cover the API and all pass."),
        ("Statistics checked", "results were compared against known cases."),
        ("Data audit", "a review of the figures found and fixed real errors."),
        ("Handles bad input", "the system stays stable with invalid or extreme data."),
    ],
    note="24 automated tests, all passing, and that number matches the report exactly. The "
         "statistics were checked against known cases, and a data audit found and fixed real "
         "errors, including a routine that hid some data and a unit that was wrong by a factor "
         "of a thousand. If asked how I know it is correct, this slide is the answer.")

# ─────────────────────────── SLIDE 18 ───────────────────────────
content("Results", "Results", 18, img=FIG / "fig4_10_heatmap.png",
    sub="Clear, verifiable outputs that meet each objective.",
    items=[
        ("122 indicators, 12,100 records", "all held in one store you can query."),
        ("15 data API endpoints", "documented and at HATEOAS Level 3."),
        ("100 out of 100", "a perfect accessibility score."),
        ("Verified figures", "every displayed number checked against the database."),
    ], caption="Figure 4.10 — The aggregation at a glance: 122 indicators × 1960–2026.",
    note="Tie each result back to an objective from slide four: the store and its coverage, the "
         "live endpoints, the accessibility score, and the audit confirming displayed figures "
         "match the stored data. All of these can be shown on the live site if asked.")

# ─────────────────────────── SLIDE 19 ───────────────────────────
content("Conclusion", "Contribution to Knowledge", 19,
    sub="A reusable result that can be rebuilt from the source code alone.",
    items=[
        ("A working reference platform", "for unified Nigerian public economic data."),
        ("A genuinely open API", "with no freely available equivalent found."),
        ("Built with free tools", "reproducible entirely from the repository."),
        ("Evidence for the argument", "that unifying this data was always achievable."),
    ],
    note="State the contribution modestly and precisely: a free, open, reproducible reference "
         "platform with a HATEOAS API, for which I found no freely available equivalent. I say "
         "no equivalent found rather than first ever, because an absolute claim is the kind a "
         "panel can pick apart.")

# ─────────────────────────── SLIDE 20 ───────────────────────────
content("Future Work", "Limitations and Future Work", 20,
    sub="Clear next steps that build on the current prototype.",
    items=[
        ("Automate collection", "connect directly to the source websites."),
        ("Expand coverage", "add more indicators, including state-level data."),
        ("Add richer analytics", "such as seasonal adjustment and forecasting."),
        ("Add API keys", "with rate limits for heavier users."),
    ],
    note="Presenting limitations as planned next steps shows I understand the boundaries rather "
         "than being caught by them: automate collection, expand coverage, add richer analytics, "
         "and add API keys with rate limits. Owning these openly is a strength, not a weakness.")

# ─────────────────────────── SLIDE 21: CLOSE ───────────────────────────
s = prs.slides.add_slide(BLANK); bg(s)
accent_bar(s, Inches(0.9), Inches(2.0), Inches(0.7))
t = box(s, Inches(0.88), Inches(2.2), Inches(11.5), Inches(1.4))
set_run(t.text_frame.paragraphs[0].add_run(), "Conclusion", 40, CREAM, HEAD_F, italic=True)
c = box(s, Inches(0.95), Inches(3.35), Inches(11.4), Inches(1.6))
cp = c.text_frame.paragraphs[0]
set_run(cp.add_run(), "Nigeria's scattered, non-machine-readable public economic data need not stay that way: it can be consolidated into a correct, programmatically usable resource using only free and open-source tools.", 18, SLATE, BODY_F)
links = box(s, Inches(0.95), Inches(5.05), Inches(11.4), Inches(1.6)); lf = links.text_frame
rows = [
    ("Dashboard:  ", "antd-cr7.github.io/nigerian-dashboard"),
    ("Open API:  ", "npedata-api.onrender.com  (docs at /docs)"),
    ("Source:  ", "github.com/ANTD-CR7/nigerian-dashboard"),
]
for i, (k, v) in enumerate(rows):
    p = lf.paragraphs[0] if i == 0 else lf.add_paragraph(); p.space_after = Pt(4)
    set_run(p.add_run(), k, 14, GREEN, BODY_F, bold=True)
    set_run(p.add_run(), v, 14, GOLD, BODY_F)
ty = box(s, Inches(0.9), Inches(6.6), Inches(11), Inches(0.7))
set_run(ty.text_frame.paragraphs[0].add_run(), "Thank you.  Questions & live demonstration.", 22, CREAM, HEAD_F, italic=True)
notes(s, "Close on the integrity line, that a platform which tells you when it might be "
         "misleading you is the real contribution. Thank the panel and invite questions and the "
         "live demonstration. Have both the live site and the offline backup open and ready.")

prs.save(OUT)
print(f"Wrote {OUT}  ({OUT.stat().st_size:,} bytes, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")
